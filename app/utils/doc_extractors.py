from docx import Document
from pptx import Presentation


def extract_text_from_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    slides_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slides_text.append(shape.text)

    return "\n".join(slides_text)